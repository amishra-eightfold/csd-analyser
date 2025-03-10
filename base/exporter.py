"""Base class for file export operations."""

import pandas as pd
from typing import Optional, Dict, Any, List, Union
from datetime import datetime
from io import BytesIO
from utils.error_handlers import handle_errors
from pptx import Presentation
from pptx.util import Inches, Pt

class BaseExporter:
    """Base class for exporting data to various formats."""
    
    def __init__(self, df: Optional[pd.DataFrame] = None):
        """
        Initialize the exporter.
        
        Args:
            df (Optional[pd.DataFrame]): DataFrame to export
        """
        self.df = df.copy() if df is not None else None
        
    def set_data(self, df: pd.DataFrame):
        """
        Set the DataFrame to export.
        
        Args:
            df (pd.DataFrame): DataFrame to export
        """
        self.df = df.copy()
    
    @handle_errors(custom_message="Error exporting to Excel")
    def to_excel(self, 
                 filename: str,
                 sheets: Optional[Dict[str, Union[pd.DataFrame, Dict[str, Any]]]] = None,
                 include_summary: bool = True) -> BytesIO:
        """
        Export data to Excel format.
        
        Args:
            filename (str): Name of the output file
            sheets (Optional[Dict[str, Union[pd.DataFrame, Dict[str, Any]]]]): Additional sheets to include
            include_summary (bool): Whether to include a summary sheet
            
        Returns:
            BytesIO: Excel file content
        """
        if self.df is None and not sheets:
            raise ValueError("No data to export")
            
        output = BytesIO()
        
        with pd.ExcelWriter(output, engine='openpyxl') as writer:
            # Write main data if available
            if self.df is not None:
                self.df.to_excel(writer, sheet_name='Data', index=False)
            
            # Write additional sheets
            if sheets:
                for sheet_name, sheet_data in sheets.items():
                    if isinstance(sheet_data, pd.DataFrame):
                        sheet_data.to_excel(writer, sheet_name=sheet_name, index=False)
                    elif isinstance(sheet_data, dict):
                        pd.DataFrame(sheet_data).to_excel(writer, sheet_name=sheet_name)
            
            # Add summary sheet if requested
            if include_summary and self.df is not None:
                summary_data = {
                    'Metric': ['Total Rows', 'Total Columns', 'Missing Values'],
                    'Value': [
                        len(self.df),
                        len(self.df.columns),
                        self.df.isnull().sum().sum()
                    ]
                }
                pd.DataFrame(summary_data).to_excel(writer, sheet_name='Summary', index=False)
        
        return output
    
    @handle_errors(custom_message="Error exporting to CSV")
    def to_csv(self, filename: str) -> BytesIO:
        """
        Export data to CSV format.
        
        Args:
            filename (str): Name of the output file
            
        Returns:
            BytesIO: CSV file content
        """
        if self.df is None:
            raise ValueError("No data to export")
            
        output = BytesIO()
        self.df.to_csv(output, index=False)
        return output
    
    @handle_errors(custom_message="Error exporting to PowerPoint")
    def to_powerpoint(self,
                     filename: str,
                     title: str = "Data Analysis Report",
                     include_charts: bool = True) -> BytesIO:
        """
        Export data to PowerPoint format.
        
        Args:
            filename (str): Name of the output file
            title (str): Title for the presentation
            include_charts (bool): Whether to include charts
            
        Returns:
            BytesIO: PowerPoint file content
        """
        if self.df is None:
            raise ValueError("No data to export")
            
        # Create presentation
        prs = Presentation()
        
        # Set slide width and height (16:9 aspect ratio)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = title
        subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        
        # Overview slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = "Data Overview"
        
        # Add overview statistics
        stats_text = (
            f"• Total Records: {len(self.df)}\n"
            f"• Total Features: {len(self.df.columns)}\n"
            f"• Missing Values: {self.df.isnull().sum().sum()}\n"
            f"• Numeric Columns: {len(self.df.select_dtypes(include=['number']).columns)}\n"
            f"• Categorical Columns: {len(self.df.select_dtypes(include=['object']).columns)}"
        )
        content.text = stats_text
        
        # Save presentation
        output = BytesIO()
        prs.save(output)
        return output
    
    @handle_errors(custom_message="Error exporting to JSON")
    def to_json(self, filename: str, orient: str = 'records') -> BytesIO:
        """
        Export data to JSON format.
        
        Args:
            filename (str): Name of the output file
            orient (str): Orientation of the JSON export
            
        Returns:
            BytesIO: JSON file content
        """
        if self.df is None:
            raise ValueError("No data to export")
            
        output = BytesIO()
        output.write(self.df.to_json(orient=orient).encode())
        return output
    
    @handle_errors(custom_message="Error exporting to HTML")
    def to_html(self, 
                filename: str,
                title: str = "Data Report",
                include_styles: bool = True) -> BytesIO:
        """
        Export data to HTML format.
        
        Args:
            filename (str): Name of the output file
            title (str): Title for the HTML document
            include_styles (bool): Whether to include CSS styles
            
        Returns:
            BytesIO: HTML file content
        """
        if self.df is None:
            raise ValueError("No data to export")
            
        # Basic CSS styles
        styles = """
        <style>
            body { font-family: Arial, sans-serif; margin: 20px; }
            table { border-collapse: collapse; width: 100%; }
            th, td { border: 1px solid #ddd; padding: 8px; text-align: left; }
            th { background-color: #f2f2f2; }
            tr:nth-child(even) { background-color: #f9f9f9; }
            h1 { color: #333; }
        </style>
        """ if include_styles else ""
        
        # Create HTML content
        html_content = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <title>{title}</title>
            {styles}
        </head>
        <body>
            <h1>{title}</h1>
            <p>Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}</p>
            {self.df.to_html(index=False)}
        </body>
        </html>
        """
        
        output = BytesIO()
        output.write(html_content.encode())
        return output 