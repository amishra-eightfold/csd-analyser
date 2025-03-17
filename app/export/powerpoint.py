"""PowerPoint export module for CSD Analyzer."""

from typing import Dict, List, Any, Optional
import pandas as pd
import numpy as np
from datetime import datetime
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from io import BytesIO
import plotly.io as pio
from ..analysis.ticket_analysis import TicketAnalyzer
from ..analysis.pattern_analysis import PatternAnalyzer
from ..visualization.charts import ChartGenerator

logger = logging.getLogger(__name__)

class PowerPointExporter:
    """Exports support data analysis to PowerPoint format."""
    
    def __init__(self, df: pd.DataFrame):
        """
        Initialize PowerPoint exporter.
        
        Args:
            df (pd.DataFrame): DataFrame containing support ticket data
        """
        self.df = df
        self.ticket_analyzer = TicketAnalyzer(df)
        self.pattern_analyzer = PatternAnalyzer(df)
        self.chart_gen = ChartGenerator(df)
        
        # Define styles
        self.title_font_size = Pt(32)
        self.heading_font_size = Pt(24)
        self.body_font_size = Pt(14)
        self.accent_color = RGBColor(31, 78, 120)  # #1F4E78
        
    def export_analysis(self,
                       filename: str,
                       title: str = "Support Analysis Report",
                       customers: Optional[List[str]] = None) -> str:
        """
        Export complete analysis to PowerPoint file.
        
        Args:
            filename (str): Output filename
            title (str): Presentation title
            customers (Optional[List[str]]): List of customers to include
            
        Returns:
            str: Path to generated PowerPoint file
        """
        try:
            # Create presentation
            prs = Presentation()
            
            # Set slide dimensions (16:9)
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            
            # Generate slides
            self._create_title_slide(prs, title)
            self._create_overview_slide(prs)
            self._create_trends_slide(prs)
            self._create_patterns_slide(prs)
            self._create_product_analysis_slide(prs)
            self._create_critical_issues_slide(prs)
            self._create_recommendations_slide(prs)
            
            # Save presentation
            prs.save(filename)
            return filename
            
        except Exception as e:
            logger.error(f"Failed to export analysis to PowerPoint: {str(e)}")
            raise
            
    def _create_title_slide(self, prs: Presentation, title: str):
        """Create title slide."""
        try:
            # Use title slide layout
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            
            # Set title
            title_shape = slide.shapes.title
            title_shape.text = title
            title_shape.text_frame.paragraphs[0].font.size = self.title_font_size
            title_shape.text_frame.paragraphs[0].font.color.rgb = self.accent_color
            
            # Add subtitle with date
            subtitle = slide.placeholders[1]
            subtitle.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
            subtitle.text_frame.paragraphs[0].font.size = self.heading_font_size
            
        except Exception as e:
            logger.error(f"Failed to create title slide: {str(e)}")
            raise
            
    def _create_overview_slide(self, prs: Presentation):
        """Create overview slide with key metrics."""
        try:
            # Use section header layout
            slide = prs.slides.add_slide(prs.slide_layouts[2])
            
            # Set title
            title_shape = slide.shapes.title
            title_shape.text = "Support Overview"
            title_shape.text_frame.paragraphs[0].font.size = self.heading_font_size
            title_shape.text_frame.paragraphs[0].font.color.rgb = self.accent_color
            
            # Get metrics
            metrics = self.ticket_analyzer.get_basic_metrics()
            
            # Create text box for metrics
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(4)
            height = Inches(4)
            
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            
            # Add metrics
            metric_text = [
                f"Total Tickets: {metrics['total_tickets']:,}",
                f"Average Resolution Time: {metrics['avg_resolution_time']:.1f} days",
                f"CSAT Score: {metrics['avg_csat']:.2f}",
                f"Escalation Rate: {metrics['escalation_rate']:.1%}"
            ]
            
            for text in metric_text:
                p = tf.add_paragraph()
                p.text = text
                p.font.size = self.body_font_size
                p.space_after = Pt(12)
            
            # Add charts
            # Volume trend
            chart_img = self._fig_to_image(
                self.chart_gen.create_ticket_volume_chart(title="")
            )
            
            left = Inches(5)
            top = Inches(1.5)
            width = Inches(7)
            height = Inches(3)
            
            slide.shapes.add_picture(
                chart_img,
                left, top, width, height
            )
            
            # Priority distribution
            chart_img = self._fig_to_image(
                self.chart_gen.create_priority_distribution_chart(title="")
            )
            
            left = Inches(5)
            top = Inches(4.5)
            width = Inches(3.5)
            height = Inches(2.5)
            
            slide.shapes.add_picture(
                chart_img,
                left, top, width, height
            )
            
            # Resolution time
            chart_img = self._fig_to_image(
                self.chart_gen.create_resolution_time_chart(title="")
            )
            
            left = Inches(8.5)
            top = Inches(4.5)
            width = Inches(3.5)
            height = Inches(2.5)
            
            slide.shapes.add_picture(
                chart_img,
                left, top, width, height
            )
            
        except Exception as e:
            logger.error(f"Failed to create overview slide: {str(e)}")
            raise
            
    def _create_trends_slide(self, prs: Presentation):
        """Create trends analysis slide."""
        try:
            # Use section header layout
            slide = prs.slides.add_slide(prs.slide_layouts[2])
            
            # Set title
            title_shape = slide.shapes.title
            title_shape.text = "Support Trends"
            title_shape.text_frame.paragraphs[0].font.size = self.heading_font_size
            title_shape.text_frame.paragraphs[0].font.color.rgb = self.accent_color
            
            # Get trends
            trends = self.pattern_analyzer.analyze_trends()
            
            # Create text box for trend summary
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(4)
            height = Inches(4)
            
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            
            # Add trend summary
            trend_text = [
                f"Volume Trend: {trends.get('volume_trend', 'N/A')}",
                f"Resolution Time Trend: {trends.get('resolution_trend', 'N/A')}",
                f"CSAT Trend: {trends.get('csat_trend', 'N/A')}"
            ]
            
            for text in trend_text:
                p = tf.add_paragraph()
                p.text = text
                p.font.size = self.body_font_size
                p.space_after = Pt(12)
            
            # Add charts
            # CSAT trend
            chart_img = self._fig_to_image(
                self.chart_gen.create_csat_trend_chart(title="")
            )
            
            left = Inches(5)
            top = Inches(1.5)
            width = Inches(7)
            height = Inches(5)
            
            slide.shapes.add_picture(
                chart_img,
                left, top, width, height
            )
            
        except Exception as e:
            logger.error(f"Failed to create trends slide: {str(e)}")
            raise
            
    def _create_patterns_slide(self, prs: Presentation):
        """Create patterns analysis slide."""
        try:
            # Use section header layout
            slide = prs.slides.add_slide(prs.slide_layouts[2])
            
            # Set title
            title_shape = slide.shapes.title
            title_shape.text = "Support Patterns"
            title_shape.text_frame.paragraphs[0].font.size = self.heading_font_size
            title_shape.text_frame.paragraphs[0].font.color.rgb = self.accent_color
            
            # Get patterns
            patterns = self.pattern_analyzer.analyze_time_patterns()
            
            # Create text box for pattern summary
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(4)
            height = Inches(4)
            
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            
            # Add daily patterns
            daily = patterns.get('daily', {})
            if daily:
                p = tf.add_paragraph()
                p.text = "Daily Patterns"
                p.font.size = self.body_font_size
                p.font.bold = True
                p.space_after = Pt(6)
                
                p = tf.add_paragraph()
                p.text = f"Peak Hours: {', '.join(map(str, daily['peak_hours']))}"
                p.font.size = self.body_font_size
                p.space_after = Pt(6)
                
                p = tf.add_paragraph()
                p.text = f"Avg Tickets/Hour: {daily['avg_tickets_per_hour']:.1f}"
                p.font.size = self.body_font_size
                p.space_after = Pt(12)
            
            # Add weekly patterns
            weekly = patterns.get('weekly', {})
            if weekly:
                p = tf.add_paragraph()
                p.text = "Weekly Patterns"
                p.font.size = self.body_font_size
                p.font.bold = True
                p.space_after = Pt(6)
                
                p = tf.add_paragraph()
                p.text = f"Busiest Days: {', '.join(weekly['busy_days'])}"
                p.font.size = self.body_font_size
                p.space_after = Pt(6)
                
                p = tf.add_paragraph()
                p.text = f"Avg Tickets/Day: {weekly['avg_tickets_per_day']:.1f}"
                p.font.size = self.body_font_size
                p.space_after = Pt(12)
            
            # Add correlation matrix
            chart_img = self._fig_to_image(
                self.chart_gen.create_correlation_matrix(title="")
            )
            
            left = Inches(5)
            top = Inches(1.5)
            width = Inches(7)
            height = Inches(5)
            
            slide.shapes.add_picture(
                chart_img,
                left, top, width, height
            )
            
        except Exception as e:
            logger.error(f"Failed to create patterns slide: {str(e)}")
            raise
            
    def _create_product_analysis_slide(self, prs: Presentation):
        """Create product area analysis slide."""
        try:
            # Use section header layout
            slide = prs.slides.add_slide(prs.slide_layouts[2])
            
            # Set title
            title_shape = slide.shapes.title
            title_shape.text = "Product Analysis"
            title_shape.text_frame.paragraphs[0].font.size = self.heading_font_size
            title_shape.text_frame.paragraphs[0].font.color.rgb = self.accent_color
            
            # Add charts
            # Ticket volume by product area
            chart_img = self._fig_to_image(
                self.chart_gen.create_product_area_chart(metric='count', title="")
            )
            
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(6)
            height = Inches(2.5)
            
            slide.shapes.add_picture(
                chart_img,
                left, top, width, height
            )
            
            # Resolution time by product area
            chart_img = self._fig_to_image(
                self.chart_gen.create_product_area_chart(metric='resolution_time', title="")
            )
            
            left = Inches(6.5)
            top = Inches(1.5)
            width = Inches(6)
            height = Inches(2.5)
            
            slide.shapes.add_picture(
                chart_img,
                left, top, width, height
            )
            
            # CSAT by product area
            chart_img = self._fig_to_image(
                self.chart_gen.create_product_area_chart(metric='csat', title="")
            )
            
            left = Inches(0.5)
            top = Inches(4.5)
            width = Inches(6)
            height = Inches(2.5)
            
            slide.shapes.add_picture(
                chart_img,
                left, top, width, height
            )
            
        except Exception as e:
            logger.error(f"Failed to create product analysis slide: {str(e)}")
            raise
            
    def _create_critical_issues_slide(self, prs: Presentation):
        """Create critical issues slide."""
        try:
            # Use section header layout
            slide = prs.slides.add_slide(prs.slide_layouts[2])
            
            # Set title
            title_shape = slide.shapes.title
            title_shape.text = "Critical Issues"
            title_shape.text_frame.paragraphs[0].font.size = self.heading_font_size
            title_shape.text_frame.paragraphs[0].font.color.rgb = self.accent_color
            
            # Get critical issues
            critical_issues = self.ticket_analyzer.get_critical_issues()
            
            if critical_issues:
                # Create table
                rows = len(critical_issues) + 1  # +1 for header
                cols = 5
                left = Inches(0.5)
                top = Inches(1.5)
                width = Inches(12)
                height = Inches(0.6 * rows)
                
                table = slide.shapes.add_table(
                    rows, cols, left, top, width, height
                ).table
                
                # Set headers
                headers = [
                    "Case Number",
                    "Subject",
                    "Priority",
                    "Status",
                    "Resolution Time"
                ]
                
                for i, header in enumerate(headers):
                    cell = table.cell(0, i)
                    cell.text = header
                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.font.size = self.body_font_size
                    paragraph.font.bold = True
                
                # Add data
                for row, issue in enumerate(critical_issues, 1):
                    values = [
                        issue['case_number'],
                        issue['subject'],
                        issue['priority'],
                        issue['status'],
                        f"{issue.get('resolution_time', 'N/A')} days"
                    ]
                    
                    for col, value in enumerate(values):
                        cell = table.cell(row, col)
                        cell.text = str(value)
                        paragraph = cell.text_frame.paragraphs[0]
                        paragraph.font.size = self.body_font_size
            else:
                # Add message if no critical issues
                left = Inches(0.5)
                top = Inches(1.5)
                width = Inches(12)
                height = Inches(1)
                
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                
                p = tf.add_paragraph()
                p.text = "No critical issues identified"
                p.font.size = self.body_font_size
                p.alignment = PP_ALIGN.CENTER
            
        except Exception as e:
            logger.error(f"Failed to create critical issues slide: {str(e)}")
            raise
            
    def _create_recommendations_slide(self, prs: Presentation):
        """Create recommendations slide."""
        try:
            # Use section header layout
            slide = prs.slides.add_slide(prs.slide_layouts[2])
            
            # Set title
            title_shape = slide.shapes.title
            title_shape.text = "Recommendations"
            title_shape.text_frame.paragraphs[0].font.size = self.heading_font_size
            title_shape.text_frame.paragraphs[0].font.color.rgb = self.accent_color
            
            # Get patterns for recommendations
            patterns = self.pattern_analyzer.identify_patterns()
            anomalies = self.pattern_analyzer.identify_anomalies()
            clusters = self.pattern_analyzer.identify_issue_clusters()
            
            # Create text box for recommendations
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(12)
            height = Inches(5)
            
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            
            # Add recommendations based on patterns
            if patterns:
                p = tf.add_paragraph()
                p.text = "Based on Pattern Analysis:"
                p.font.size = self.body_font_size
                p.font.bold = True
                p.space_after = Pt(6)
                
                for area, metrics in patterns.get('correlations', {}).items():
                    p = tf.add_paragraph()
                    p.text = f"• Review {area} performance metrics"
                    p.font.size = self.body_font_size
                    p.space_after = Pt(6)
            
            # Add recommendations based on anomalies
            if anomalies:
                p = tf.add_paragraph()
                p.text = "\nBased on Anomaly Detection:"
                p.font.size = self.body_font_size
                p.font.bold = True
                p.space_after = Pt(6)
                
                for anomaly in anomalies:
                    p = tf.add_paragraph()
                    if anomaly['type'] == 'resolution_time':
                        p.text = f"• Investigate case {anomaly['case_number']} for resolution delays"
                    else:
                        p.text = f"• Review ticket volume spike on {anomaly['date'].strftime('%Y-%m-%d')}"
                    p.font.size = self.body_font_size
                    p.space_after = Pt(6)
            
            # Add recommendations based on clusters
            if clusters:
                p = tf.add_paragraph()
                p.text = "\nBased on Issue Clusters:"
                p.font.size = self.body_font_size
                p.font.bold = True
                p.space_after = Pt(6)
                
                for cluster in clusters:
                    p = tf.add_paragraph()
                    p.text = f"• Address common issues around: {', '.join(cluster['keywords'])}"
                    p.font.size = self.body_font_size
                    p.space_after = Pt(6)
            
        except Exception as e:
            logger.error(f"Failed to create recommendations slide: {str(e)}")
            raise
            
    def _fig_to_image(self, fig: Any) -> BytesIO:
        """Convert Plotly figure to image."""
        try:
            img_bytes = pio.to_image(fig, format='png')
            return BytesIO(img_bytes)
            
        except Exception as e:
            logger.error(f"Failed to convert figure to image: {str(e)}")
            raise
